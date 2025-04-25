import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from db import fetch_data_from_sp
from io import BytesIO

def add_table_to_slide(slide, df, title):
    """Adds a formatted table to a PowerPoint slide."""
    if df.empty:
        print(f"⚠️ No data available for {title}")
        return None

    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Define table position and size
    rows, cols = df.shape[0] + 1, df.shape[1]
    x, y, width, height = Inches(0.5), Inches(1), Inches(9), Inches(1.5 + len(df) * 0.2)

    # Add table to slide
    table = slide.shapes.add_table(rows, cols, x, y, width, height).table

    # Set column headers
    for col_idx, column_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = column_name
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Inches(0.2)

    # Fill table with data
    for row_idx, row in enumerate(df.itertuples(index=False), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)

    return table

def add_chart_to_slide(slide, df, x_column, y_column, title, xlabel, ylabel):
    """Generates and inserts a bar chart into a PowerPoint slide."""
    if df.empty:
        print(f"⚠️ No data available for {title}")
        return None

    fig, ax = plt.subplots(figsize=(6, 4))
    bars = ax.bar(df[x_column], df[y_column], color='skyblue')

    # Add data labels on top of each bar
    for bar in bars:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2, height, f'{height:,}', 
                ha='center', va='bottom', fontsize=10, fontweight='bold')

    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    ax.set_title(title)
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()

    # Save chart to memory
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=300)
    plt.close()

    # Insert chart into PowerPoint
    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    slide.shapes.add_picture(img_stream, left, top, width, height)

def generate_ppt(master_client_id, booking_month, client_id, output_filename="booking_report.pptx"):
    """Fetches stored procedure data and generates PowerPoint with tables and charts."""
    dfs = fetch_data_from_sp(master_client_id, booking_month, client_id)

    if not dfs or len(dfs) < 5:
        print("⚠️ No sufficient data fetched from stored procedure.")
        return None

    # Create a PowerPoint presentation
    prs = Presentation()

    # 📌 Table 1: Client-wise Booking Statistics
    client_df = dfs[0][["ClientName", "total_guests", "total_bookings", "total_spent"]]
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    add_table_to_slide(slide, client_df, "Client Booking Statistics")

    # 📌 Table 2: Traveller Profile
    tp_df = dfs[1][["ClientName", "guest_volume_below_2000", "guest_volume_2000_5000", "guest_volume_above_5000"]]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_table_to_slide(slide, tp_df, "Traveller Profile")

    # 📌 Table 3: City-wise Booking Statistics
    city_df = dfs[3][["city", "total_guests", "total_bookings", "total_spent", "total_roomnights"]]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_table_to_slide(slide, city_df, "City Booking Statistics")

    # 📌 Table 4: Hotel-wise Booking Statistics
    hotel_df = dfs[4][["PropertyName", "total_guests", "total_bookings", "total_spent", "total_roomnights"]]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_table_to_slide(slide, hotel_df, "Hotel Booking Statistics")

    # 📊 Chart: Booking Lead Time Distribution (if available)
    if len(dfs) > 2:
        lead_time_df = dfs[2]  # Modify index if needed
        if not lead_time_df.empty:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            add_chart_to_slide(slide, lead_time_df, "LeadTimeRange", "BookingCount", 
                               "Booking Lead Time Distribution", "Lead Time Range", "Booking Count")

    # Save PowerPoint file
    prs.save(output_filename)
    print(f"✅ PowerPoint report saved as {output_filename}")

    return output_filename

if __name__ == "__main__":
    master_client_id = 969  # Replace with valid ID
    client_id = 2631  # Replace with valid Client ID for filtering
    booking_month = 'Sep 23'
    ppt_file = generate_ppt(master_client_id, booking_month, client_id)
    print(f"Generated PowerPoint file: {ppt_file}")
