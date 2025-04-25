import os
from pptx import Presentation
from pptx.util import Inches
from table import generate_tables
from chart_generator import generate_charts

def add_slide(prs, title):
    """Adds a slide with a given title."""
    slide_layout = prs.slide_layouts[5]  # Title Only Layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    return slide

def add_image_to_slide(slide, image_path):
    """Adds an image to a slide if the file exists."""
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(6))
    else:
        print(f"⚠️ Image not found: {image_path}")

def generate_ppt(master_client_id, client_id):
    """Generates a PowerPoint report with tables and charts."""
    prs = Presentation()
    
    # Fetch table images
    table_files = generate_tables(master_client_id, client_id)
    chart_files = generate_charts(master_client_id, client_id)
    
    # Add Table Slides
    for title, file in table_files.items():
        if file:
            slide = add_slide(prs, title.replace("_", " ").title())
            add_image_to_slide(slide, file)
    
    # Add Chart Slides
    for file in chart_files:
        if file:
            slide = add_slide(prs, os.path.basename(file).replace("_", " ").title())
            add_image_to_slide(slide, file)
    
    # Save Presentation
    ppt_filename = f"Booking_Report_{master_client_id}.pptx"
    prs.save(ppt_filename)
    print(f"✅ PowerPoint Report saved as {ppt_filename}")
    return ppt_filename

# 🔹 Test Function
if __name__ == "__main__":
    master_client_id = 969  # Replace with valid ID
    client_id = 803  # Replace with valid Client ID
    generate_ppt(master_client_id, client_id)
