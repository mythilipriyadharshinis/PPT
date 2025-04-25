import pandas as pd
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor  # For colors
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN  # For text alignment
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from db import fetch_data_from_sp
from io import BytesIO




def fill_thank_you_slide(slide, username):
    """Adds login link and credentials to the Thank You slide with a working hyperlink."""
    login_url = "https://powerbireport.hummingbirdindia.com"  # Replace this with your actual URL

    # Create a textbox
    left, top, width, height = Inches(0.5), Inches(9.3), Inches(10), Inches(2.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # Add hyperlink paragraph
    p1 = text_frame.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "🔗 Click here to login to the Interactive Dashboard"
    run1.hyperlink.address = login_url
    run1.font.size = Pt(24)
    run1.font.name = "DM Sans"
    run1.font.color.rgb = RGBColor(228, 33, 39)  # Blue link color

    # Add new paragraph for credentials
    p2 = text_frame.add_paragraph()
    p2.text = (
        f"Username: {username}\n"
        "Password: 123123"
    )
    p2.font.size = Pt(24)
    p2.font.name = "DM Sans"
    p2.font.color.rgb = RGBColor(0, 0, 0)



def fill_client_name(slide, client_name):
    """Fills the client name placeholder in the first slide with specific font style, size, and color."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if "{Client_name}" in para.text:
                    para.text = para.text.replace("{Client_name}", client_name)

                    # Apply font styling
                    for run in para.runs:
                        run.font.bold = True  # Make text bold
                        run.font.size = Pt(72)  # Set font size
                        run.font.name = "DM Sans"  # Set font type
                        run.font.color.rgb = RGBColor(0,0,0)  # Set text color to black

def add_table_to_slide(slide, df, max_table_height=7):
    """Adds a table to an existing slide."""
    if df.empty:
        print("⚠️ No data available.")
        return None

    # Define table position
    rows, cols = df.shape[0] + 1, df.shape[1]
    x, y, width = Inches(1.5), Inches(2), Inches(17.5)

    # Dynamically adjust row height & font size
    num_rows = df.shape[0] + 1
    base_row_height = 0.4 if num_rows > 20 else 0.6  # Shrink if many rows
    total_height = base_row_height * num_rows
    height = min(Inches(total_height), Inches(max_table_height))

    # Add table
    table = slide.shapes.add_table(rows, cols, x, y, width, height).table

    # Determine font size based on row count
    font_size_header = Pt(18 if num_rows <= 15 else 14)
    font_size_body = Pt(16 if num_rows <= 15 else 12)

    # Style Header
    for col_idx, column_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = column_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(252,156,20)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.bold = True
        para.font.size = font_size_header
        para.font.color.rgb = RGBColor(255, 255, 255)

    # Style Body
    for row_idx, row in enumerate(df.itertuples(index=False), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.size = font_size_body
            para.font.color.rgb = RGBColor(0, 0, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 230, 220)

def add_chart_to_slide(slide, df, x_column, y_column, title):
    """Adds a pie chart with data labels and a dynamically generated text summary on the right."""
    if df.empty:
        print(f"⚠️ No data available for chart: {title}")
        return None

    # Define custom colors
    custom_colors = ['#d5720a','#F5F1EB', '#D9D2C1', '#A89E8F', '#736B5E', '#4C463C', '#8A2BE2', '#FF6347', '#20B2AA']


    fig, ax = plt.subplots(figsize=(14, 12))  # Large figure
    fig.patch.set_facecolor('#F0F0F0')  # Light gray background
    wedges, texts, autotexts = ax.pie(
        df[y_column], labels=df[x_column], autopct='%1.1f%%', colors=custom_colors[:len(df)],
        startangle=140, wedgeprops={'edgecolor': 'black', 'linewidth': 1}
    )

    ax.set_title(title, fontsize=20, fontweight='bold')

    # Increase text size
    plt.setp(texts, fontsize=20)
    plt.setp(autotexts, fontsize=20, fontweight="bold")

    plt.tight_layout()

    # Save the chart as an image
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=300)
    plt.close()

    # Insert chart into slide
    left, top, width, height = Inches(2), Inches(2.5), Inches(7), Inches(6)  # Space for explanation
    slide.shapes.add_picture(img_stream, left, top, width, height)

    # **Generate dynamic insights**  
    total = df[y_column].sum()
    df['Percentage'] = (df[y_column] / total) * 100
    top_category = df.loc[df[y_column].idxmax(), x_column]
    top_percentage = df.loc[df[y_column].idxmax(), 'Percentage']

    # **Create text similar to the provided image**  
    insight_text = (
        f"🔹 {top_percentage:.1f}% of the bookings are made in {top_category}.\n\n"
        f"More structured booking planning can optimize costs, ensure availability, and enhance the overall booking experience."
    )

    # Add text box for insights
    text_left, text_top, text_width, text_height = Inches(12), Inches(3), Inches(4), Inches(3)  # Right side
    textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # Add insight text
    p = text_frame.add_paragraph()
    p.text = insight_text
    p.font.size = Pt(24)

def generate_ppt(master_client_id, booking_month, client_id=None, is_client_level=False):
    """Fills an existing 8-slide PowerPoint template with data."""
    
    TEMPLATE_PATH = "client_template.pptx" if is_client_level else "masterclient_template.pptx"
    
    #dfs = fetch_data_from_sp(master_client_id, booking_month,client_id=None, is_client_level=False)
    
    if is_client_level:
        dfs = fetch_data_from_sp(master_client_id, client_id, booking_month, is_client_level=True)
    else:
        dfs = fetch_data_from_sp(master_client_id, booking_month, client_id=None, is_client_level=False)

    if not dfs or (len(dfs) < 10 and not is_client_level) or (len(dfs) < 7 and is_client_level):
        print("⚠️ Not enough data to fill all slides.")
        return None

    prs = Presentation(TEMPLATE_PATH)

    # Step 1: Fill Client Name on Slide 1
    fill_client_name(prs.slides[0], dfs[0]["clientname"].iloc[0])
    client_name = dfs[0]["clientname"].iloc[0]
    
    if is_client_level==False:
     for col in dfs[9].columns:
      if "Reservation Value" in col:
        dfs[9][col] = dfs[9][col].apply(lambda x: f'₹ {int(x)}' if pd.notnull(x) else '₹ 0')
     
        
    # Step 2: Insert Data into Slides 2-7
    if  is_client_level == False:
      slide_data_mapping = {
        1: dfs[1],  # Slide 2 → Client Booking Statistics
        2: dfs[2],  # Slide 3 → Top cities with ADR(entity wise)
        3: dfs[3],  # Slide 4 → Top properties with ADR(entity wise)
        4: dfs[4],  # Slide 5 → Customer feedback
        5: dfs[5],  # Slide 6 → Traveller Profile
        6: dfs[6],  # Slide 7 → Lead Time Range Analysis
        7: dfs[7],  # Slide 8 → Top 10 Cities by bookings
        8: dfs[8],  # slide 9 → Top 10 Properties by bookings
        9: dfs[9],  # slide 10→ Reservation Summary -last 3 months 
    }
    else:
        slide_data_mapping = {
        1: dfs[1],  # Slide 2 → Client Booking Statistics
        2: dfs[2],  # Slide 3 →  Customer feedback
        3: dfs[3],  # Slide 4 → Traveller Profile 
        4: dfs[4],  # Slide 5 → Lead Time Range Analysis
        5: dfs[5],  # Slide 6 → Top 10 Cities by bookings
        6: dfs[6],  # Slide 7 → Top 10 Properties by bookings
        7: dfs[7],  # Slide 8 → Top 10 Grade
    }   

    for slide_idx, df in slide_data_mapping.items():
        slide = prs.slides[slide_idx]

        if "LeadTimeRange" in df.columns and "BookingCount" in df.columns:
            add_chart_to_slide(slide, df, "LeadTimeRange", "BookingCount", "Booking Lead Time")
        else:
            add_table_to_slide(slide, df)
    # Step 0: Create reports folder if it doesn't exist
    
    # Step 3: Keep Slide 9 ("Thank You") Unchanged
    # Step 4: Fill the Thank You slide (Slide 10) with login info
    if is_client_level==False:
     username = dfs[10]["UserName"].iloc[0]  # Assuming the column is named 'username'
     fill_thank_you_slide(prs.slides[10], username)
    else:
     username = dfs[8]["UserName"].iloc[0]  # Assuming the column is named 'username'
     fill_thank_you_slide(prs.slides[8], username)
    
    output_folder = "reports"
    os.makedirs(output_folder, exist_ok=True)
    
    # Step 2: Clean client name for filename
    safe_client_name = client_name.replace(" ", "_").replace("/", "_")
    
    if client_id:
     filename = f"booking_report_{safe_client_name}_{booking_month}_{client_id}.pptx"
    else:
     filename = f"booking_report_{safe_client_name}_{booking_month}_{master_client_id}.pptx"
     
     # Step 4: Full path with folder
    output_filename = os.path.join(output_folder, filename)

    # Save PowerPoint
    prs.save(output_filename)
    print(f"✅ PowerPoint report saved as {output_filename}")
    return output_filename

if __name__ == "__main__":
    generate_ppt(969,'2024-03',2631,is_client_level=True)
